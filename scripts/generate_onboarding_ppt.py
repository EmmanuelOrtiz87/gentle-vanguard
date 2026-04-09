from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(20, 40, 80)
DARK_GRAY = RGBColor(45, 45, 50)
MEDIUM_GRAY = RGBColor(80, 80, 90)
LIGHT_GRAY = RGBColor(240, 242, 245)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(0, 140, 100)
ACCENT_BLUE = RGBColor(0, 120, 200)
ORANGE = RGBColor(220, 100, 40)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf = sub.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(180, 190, 200)
        p.font.name = "Calibri"
    
    return slide

def add_header_slide(title):
    slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    return slide

def add_content_slide(title, points):
    slide = add_header_slide(title)
    
    y = 1.5
    for point in points:
        if isinstance(point, tuple):
            main, sub = point
            # Main point
            main_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.5), Inches(0.6))
            tf = main_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = main
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Calibri"
            y += 0.5
            
            # Sub point
            if sub:
                sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11), Inches(0.5))
                tf = sub_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sub
                p.font.size = Pt(16)
                p.font.color.rgb = MEDIUM_GRAY
                p.font.name = "Calibri"
                y += 0.5
            y += 0.3
        else:
            # Simple point with bullet
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.12), Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT_BLUE
            dot.line.fill.background()
            
            point_box = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(11.5), Inches(0.5))
            tf = point_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Calibri"
            y += 0.6
    
    return slide

def add_two_column_slide(title, left_title, left_points, right_title, right_points):
    slide = add_header_slide(title)
    
    # Left column
    left_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = left_label.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Calibri"
    
    y = 1.9
    for point in left_points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y + 0.12), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        
        point_box = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(5.4), Inches(0.5))
        tf = point_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        y += 0.55
    
    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.4), Inches(0.02), Inches(5.5))
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(220, 220, 230)
    divider.line.fill.background()
    
    # Right column
    right_label = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = right_label.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Calibri"
    
    y = 1.9
    for point in right_points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.9), Inches(y + 0.12), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        
        point_box = slide.shapes.add_textbox(Inches(7.2), Inches(y), Inches(5.4), Inches(0.5))
        tf = point_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        y += 0.55
    
    return slide

def add_kpi_slide(title, kpis):
    slide = add_header_slide(title)
    
    num = len(kpis)
    box_w = 2.8
    spacing = 0.3
    total_w = num * box_w + (num - 1) * spacing
    start_x = (13.333 - total_w) / 2
    
    for i, (num_str, label) in enumerate(kpis):
        x = start_x + i * (box_w + spacing)
        
        kpi_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(box_w), Inches(2.5))
        kpi_box.fill.solid()
        kpi_box.fill.fore_color.rgb = ACCENT_BLUE if i == 0 else LIGHT_GRAY
        kpi_box.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(box_w), Inches(1.2))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num_str
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE if i == 0 else NAVY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.5), Inches(box_w - 0.2), Inches(0.8))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE if i == 0 else MEDIUM_GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_code_slide(title, code_text, language=""):
    slide = add_header_slide(title)
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(30, 30, 40)
    code_box.line.fill.background()
    
    code = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(5.2))
    tf = code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(180, 200, 220)
    
    if language:
        lang_box = slide.shapes.add_textbox(Inches(11.5), Inches(1.6), Inches(1.5), Inches(0.4))
        tf = lang_box.text_frame
        p = tf.paragraphs[0]
        p.text = language.upper()
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 120, 140)
    
    return slide

def add_flow_slide(title, steps):
    slide = add_header_slide(title)
    
    y = 1.6
    for i, (step, desc) in enumerate(steps):
        # Number circle
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y), Inches(0.5), Inches(0.5))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = ACCENT_BLUE
        num_circle.line.fill.background()
        
        num_text = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.1), Inches(0.5), Inches(0.4))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        
        # Step title
        step_box = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.05), Inches(4), Inches(0.4))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = "Calibri"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.4), Inches(11), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = "Calibri"
        
        y += 0.9
    
    return slide

def add_comparison_slide(title, items):
    slide = add_header_slide(title)
    
    y = 1.6
    for term, definition in items:
        # Term box
        term_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(2.5), Inches(0.6))
        term_box.fill.solid()
        term_box.fill.fore_color.rgb = ACCENT_BLUE
        term_box.line.fill.background()
        
        term_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(2.5), Inches(0.4))
        tf = term_text.text_frame
        p = tf.paragraphs[0]
        p.text = term
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        
        # Definition
        def_box = slide.shapes.add_textbox(Inches(3.2), Inches(y + 0.1), Inches(9.5), Inches(0.5))
        tf = def_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = definition
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        
        y += 0.75
    
    return slide

def add_quickref_slide(title, commands):
    slide = add_header_slide(title)
    
    y = 1.5
    for cmd, desc in commands:
        # Command box
        cmd_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(4.5), Inches(0.5))
        cmd_box.fill.solid()
        cmd_box.fill.fore_color.rgb = RGBColor(40, 40, 50)
        cmd_box.line.fill.background()
        
        cmd_text = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.1), Inches(4.1), Inches(0.4))
        tf = cmd_text.text_frame
        p = tf.paragraphs[0]
        p.text = cmd
        p.font.size = Pt(12)
        p.font.name = "Courier New"
        p.font.color.rgb = RGBColor(100, 220, 150)
        
        # Arrow
        arrow = slide.shapes.add_textbox(Inches(5.2), Inches(y + 0.05), Inches(0.5), Inches(0.4))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_BLUE
        
        # Description
        desc_text = slide.shapes.add_textbox(Inches(5.8), Inches(y + 0.05), Inches(7), Inches(0.4))
        tf = desc_text.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        
        y += 0.65
    
    return slide

# =====================================================
# SLIDES
# =====================================================

# Slide 1: Title
add_title_slide(
    "Workspace Foundation",
    "Technical Onboarding - Guía para el Equipo de Desarrollo"
)

# Slide 2: What is it?
add_content_slide(
    "¿Qué es Workspace Foundation?",
    [
        "Sistema de plantillas y automatización para proyectos de software",
        "Integra herramientas de IA en el flujo de trabajo diario",
        "Estandariza configuración, desarrollo y documentación",
        ("Resultado:", "Más tiempo en código, menos en setup"),
        ("Beneficio principal:", "Un comando configura todo automáticamente")
    ]
)

# Slide 3: Components
add_kpi_slide(
    "Componentes Core",
    [
        ("Bootstrap", "Auto-config"),
        ("Templates", "Estructuras"),
        ("AI Tools", "Claude, OpenCode"),
        ("GGA", "Code Review")
    ]
)

# Slide 4: AI Concepts
add_content_slide(
    "Conceptos Fundamentales de IA",
    [
        ("LLM (Large Language Model):", "Modelo entrenado con texto a gran escala - GPT-4, Claude"),
        ("Prompt:", "Instrucción que le das a la IA - Cuanto más contexto, mejor respuesta"),
        ("Token:", "Unidad básica de texto que la IA procesa - 1 token ≈ 4 caracteres"),
        ("Context Window:", "Memoria máxima de la IA - 200K tokens en Claude Sonnet 4"),
        ("Temperature:", "Control de creatividad - 0 = preciso, 1 = creativo")
    ]
)

# Slide 5: Tools Overview
add_kpi_slide(
    "Herramientas Integradas",
    [
        ("Claude", "Asistencia general"),
        ("OpenCode", "CLI multi-model"),
        ("Gentle-AI", "Contextual"),
        ("GGA", "Code Review")
    ]
)

# Slide 6: Claude Code
add_two_column_slide(
    "Claude Code",
    "Casos de Uso",
    [
        "Generar código nuevo",
        "Explicar código existente",
        "Debugging y errores",
        "Refactoring",
        "Crear tests unitarios"
    ],
    "Ejemplos",
    [
        '/generate "Crea una función que..."',
        '/review "Explica este error..."',
        '/refactor "Mejora esta función..."',
        '/test "Genera tests para..."',
        '/explain "Qué hace este código?"'
    ]
)

# Slide 7: OpenCode
add_two_column_slide(
    "OpenCode",
    "Características",
    [
        "Interface CLI unificada",
        "Múltiples modelos (Claude, GPT-4)",
        "Configuración por proyecto",
        "Historial de conversaciones",
        "Integración con Git"
    ],
    "Comandos",
    [
        'opencode --model gpt-4o "request"',
        'opencode --file src/app.ts "mejora"',
        'opencode config set default-model claude',
        'opencode --help'
    ]
)

# Slide 8: GGA
add_content_slide(
    "GGA - Gentleman Guardian Angel",
    [
        "Sistema de hooks de pre-commit que revisa código con IA",
        ("Workflow:", "Git commit → GGA hook → Revisión IA → Aprueba/Bloquea"),
        "Aplica estándares automáticamente",
        "Feedback inmediato en cada commit",
        ("Para forzar commit (usar con precaución):", "git commit --no-verify -m 'mensaje'")
    ]
)

# Slide 9: Daily Workflow
add_flow_slide(
    "Flujo de Trabajo Diario",
    [
        ("Mañana - Iniciar:", "./scripts/init-workspace.ps1"),
        ("Trabajar:", "Usar AI tools + Git normalmente"),
        ("Revisar:", "GGA corre automáticamente en pre-commit"),
        ("Fin de día:", "./scripts/finalize-session.ps1"),
        ("Resultado:", "Métricas generadas automáticamente")
    ]
)

# Slide 10: Commands
add_quickref_slide(
    "Comandos Esenciales",
    [
        ("./scripts/init-workspace.ps1", "Iniciar día con bootstrap"),
        ("claude \"tu pregunta\"", "Usar Claude Code"),
        ("opencode --model claude \"request\"", "Usar OpenCode"),
        ("gentle-ai \"consulta\"", "Usar Gentle-AI"),
        ("./scripts/finalize-session.ps1", "Finalizar día con métricas")
    ]
)

# Slide 11: Prompt Best Practices
add_two_column_slide(
    "Prompt Engineering - Mejores Prácticas",
    "Hacer",
    [
        "Ser específico y detallado",
        "Incluir contexto del proyecto",
        "Especificar lenguaje/framework",
        "Pedir ejemplos del output",
        "Iterar y refinar"
    ],
    "Evitar",
    [
        "Pedidos vagos: 'haz algo'",
        "Sin contexto",
        "Pedir múltiples cosas a la vez",
        "No verificar el output",
        "Enviar información sensible"
    ]
)

# Slide 12: Security
add_content_slide(
    "Seguridad - NO Enviar a IA",
    [
        "Credenciales y API keys",
        "Contraseñas y tokens",
        "Datos de usuarios/customers",
        "Secrets de producción",
        "Información personal identificable (PII)",
        ("Lo que SÍ puedes enviar:", "Código genérico, arquitectura, patrones, errores públicos")
    ]
)

# Slide 13: Quick Reference
add_quickref_slide(
    "Quick Reference",
    [
        ("Generar código:", 'claude "Crea una función que..."'),
        ("Debug:", 'claude "Explica este error: [error]"'),
        ("Tests:", 'claude "Genera tests con Jest para..."'),
        ("Refactor:", 'claude "Refactoriza esta función para..."'),
        ("Review:", 'claude "Revisa este código"')
    ]
)

# Slide 14: Glossary 1
add_comparison_slide(
    "Glosario - Términos de IA",
    [
        ("AI", "Artificial Intelligence - Inteligencia artificial"),
        ("LLM", "Large Language Model - Modelo de lenguaje grande"),
        ("ML", "Machine Learning - Sistemas que aprenden de datos"),
        ("NLP", "Natural Language Processing - Procesamiento de lenguaje"),
        ("Prompt", "Instrucción dada a una IA"),
        ("Token", "Unidad básica de texto procesada por IA")
    ]
)

# Slide 15: Glossary 2
add_comparison_slide(
    "Glosario - Términos del Proyecto",
    [
        ("Bootstrap", "Script que configura el entorno automáticamente"),
        ("Template", "Estructura base para nuevos proyectos"),
        ("Skill", "Patrón reutilizable de prompts/acciones"),
        ("Hook", "Script que corre en eventos de Git"),
        ("Agent", "Instancia de AI tool en una máquina"),
        ("Audit", "Registro de actividad para métricas")
    ]
)

# Slide 16: FAQ
add_two_column_slide(
    "FAQ - Preguntas Frecuentes",
    "¿La IA puede ver todos mis archivos?",
    [
        "Sí cuando le das contexto.",
        "Audit system registra qué se accede.",
        "No envía código excepto vía APIs configuradas."
    ],
    "¿Qué pasa si la IA da código incorrecto?",
    [
        "GGA y code review son tu red de seguridad.",
        "IA es asistencia, no reemplazo.",
        "Siempre verificar antes de aplicar."
    ]
)

# Slide 17: FAQ 2
add_two_column_slide(
    "FAQ - Continúa",
    "¿Mis API keys están seguras?",
    [
        "Keys en .env (gitignored).",
        "Nunca compartidas con terceros.",
        "Almacenamiento local."
    ],
    "¿Puedo usar múltiples AI tools?",
    [
        "Sí, cada una tiene fortalezas.",
        "Claude para código general.",
        "GGA para review automático."
    ]
)

# Slide 18: Resources
add_content_slide(
    "Recursos y Documentación",
    [
        "Repo principal: github.com/EmmanuelOrtiz87/AI-development-stack",
        "docs/TECHNICAL-ONBOARDING.md - Guía completa",
        "docs/audit-system.md - Sistema de métricas",
        "AGENTS.md - Reglas para AI agents",
        "docs/ - Toda la documentación"
    ]
)

# Slide 19: Next Steps
add_flow_slide(
    "Próximos Pasos",
    [
        ("1. Instalar:", "./scripts/init-workspace.ps1"),
        ("2. Configurar:", "Seguir guía en docs/"),
        ("3. Practicar:", "Generar código, tests, refactor"),
        ("4. Integrar:", "Usar en tu flujo diario"),
        ("5. Medir:", "Ver métricas en .audit/")
    ]
)

# Slide 20: Questions
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

title = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "¿Preguntas?"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
tf = sub.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Workspace Foundation\nTechnical Onboarding"
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(180, 190, 200)
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "C:\\Workspace_local\\workspace-foundation\\docs\\WORKSPACE-FOUNDATION-ONBOARDING.pptx"
prs.save(output_path)
print(f"Presentación guardada: {output_path}")
