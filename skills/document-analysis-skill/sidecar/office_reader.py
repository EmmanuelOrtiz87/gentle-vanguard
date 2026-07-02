import os, json, csv as csv_lib, io, re, xml.etree.ElementTree as ET
from typing import Any

try:
    from docx import Document as DocxDocument; HAS_DOCX = True
except ImportError: HAS_DOCX = False
try:
    import openpyxl; HAS_XLSX = True
except ImportError: HAS_XLSX = False
try:
    from pptx import Presentation; HAS_PPTX = True
except ImportError: HAS_PPTX = False
try:
    import fitz; HAS_FITZ = True
except ImportError: HAS_FITZ = False
try:
    import markdown as md_lib; import yaml; HAS_MD = True
except ImportError: HAS_MD = False
try:
    from PIL import Image; HAS_PIL = True
except ImportError: HAS_PIL = False

def _safe_read(path, mode="rb"):
    with open(path, mode) as f: return f.read()

def _detect_encoding(path):
    for enc in ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]:
        try:
            with open(path, "r", encoding=enc) as f: f.read(); return enc
        except (UnicodeDecodeError, UnicodeError): continue
    return "utf-8"

def _parse_frontmatter(text):
    fm, body = {}, text
    if text.startswith("---"):
        parts = text.split("---", 2)
        if len(parts) >= 3:
            try: fm = yaml.safe_load(parts[1]) or {}
            except Exception: fm = {}
            body = parts[2].strip()
    return fm, body

def read_docx(path):
    if not HAS_DOCX: return _fallback_text(path, "docx")
    doc = DocxDocument(path)
    paragraphs, tables, images, headers, footers = [], [], [], [], []
    for p in doc.paragraphs:
        paragraphs.append({"text": p.text, "style": p.style.name if p.style else None, "alignment": str(p.alignment) if p.alignment else None})
    for i, table in enumerate(doc.tables):
        tables.append({"index": i, "rows": [[cell.text for cell in row.cells] for row in table.rows]})
    for s in doc.sections:
        if s.header: headers.append(" ".join(p.text for p in s.header.paragraphs))
        if s.footer: footers.append(" ".join(p.text for p in s.footer.paragraphs))
    return {"content": "\n".join(p["text"] for p in paragraphs),
        "metadata": {"title": doc.core_properties.title or "", "author": doc.core_properties.author or "",
            "created": str(doc.core_properties.created) if doc.core_properties.created else "",
            "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
            "page_count": len(doc.sections), "paragraph_count": len(paragraphs), "table_count": len(tables)},
        "structure": {"paragraphs": paragraphs}, "tables": tables, "images": images, "headers": headers, "footers": footers, "errors": []}

def read_xlsx(path):
    if not HAS_XLSX: return _fallback_text(path, "xlsx")
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    sheets, tables = [], []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data = [list(row) for row in ws.iter_rows(values_only=True)]
        sheets.append({"name": sheet_name, "rows": rows_data, "row_count": len(rows_data)})
        tables.append({"name": sheet_name, "rows": rows_data})
    metadata = {"sheet_count": len(wb.sheetnames), "sheet_names": wb.sheetnames}
    wb.close()
    return {"content": sheets, "metadata": metadata, "structure": {"sheets": sheets}, "tables": tables, "images": [], "errors": []}

def read_pptx(path):
    if not HAS_PPTX: return _fallback_text(path, "pptx")
    prs = Presentation(path)
    slides, images = [], []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts, slide_images = [], []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip(): slide_texts.append(shape.text)
            if shape.shape_type == 13: slide_images.append({"name": shape.name, "left": str(shape.left), "top": str(shape.top)})
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame else ""
        slides.append({"slide_number": slide_num, "texts": slide_texts, "notes": notes, "layout": slide.slide_layout.name if slide.slide_layout else "", "image_count": len(slide_images)})
        images.extend(slide_images)
    content = "\n".join(" ".join(s["texts"]) for s in slides)
    return {"content": content, "metadata": {"slide_count": len(prs.slides), "slide_width": str(prs.slide_width), "slide_height": str(prs.slide_height)},
        "structure": {"slides": slides}, "tables": [], "images": images, "errors": []}

def read_pdf(path):
    if not HAS_FITZ: return _fallback_text(path, "pdf")
    doc = fitz.open(path)
    pages, images, outlines, text_content = [], [], [], []
    for item in doc.get_toc(): outlines.append({"level": item[0], "title": item[1], "page": item[2]})
    for page_num in range(len(doc)):
        page = doc[page_num]; text = page.get_text("text"); text_content.append(text)
        for img in page.get_images(full=True):
            bi = doc.extract_image(img[0])
            images.append({"xref": img[0], "width": bi.get("width"), "height": bi.get("height"), "ext": bi.get("ext")})
        pages.append({"page_number": page_num + 1, "text": text, "char_count": len(text), "image_count": len(page.get_images(full=True))})
    meta = doc.metadata or {}; doc.close()
    return {"content": "\n".join(text_content),
        "metadata": {"title": meta.get("title", ""), "author": meta.get("author", ""), "subject": meta.get("subject", ""),
            "keywords": meta.get("keywords", ""), "producer": meta.get("producer", ""), "creator": meta.get("creator", ""),
            "page_count": len(pages), "image_count": len(images)},
        "structure": {"pages": pages, "outlines": outlines}, "tables": [], "images": images, "errors": []}

def read_txt(path):
    encoding = _detect_encoding(path)
    try:
        with open(path, "r", encoding=encoding) as f: content = f.read()
    except Exception: content = _safe_read(path, "rb").decode("utf-8", errors="replace")
    lines = content.splitlines()
    return {"content": content, "metadata": {"encoding": encoding, "line_count": len(lines), "char_count": len(content)},
        "structure": {"lines": lines}, "tables": [], "images": [], "errors": []}

def read_csv(path, delimiter=None):
    content = _safe_read(path).decode("utf-8-sig", errors="replace")
    sniffer = csv_lib.Sniffer()
    try: delimiter = sniffer.sniff(content[:4096]).delimiter if delimiter is None else delimiter
    except Exception: delimiter = delimiter or ","
    reader = csv_lib.reader(io.StringIO(content), delimiter=delimiter)
    rows = [row for row in reader]
    headers = rows[0] if rows else []
    return {"content": "\n".join(",".join(r) for r in rows),
        "metadata": {"delimiter": delimiter, "row_count": len(rows), "column_count": len(headers) if headers else 0, "headers": headers},
        "structure": {"headers": headers, "rows": rows}, "tables": [{"name": "csv_data", "rows": rows}], "images": [], "errors": []}

def read_markdown(path):
    text = _safe_read(path).decode("utf-8", errors="replace")
    fm, body = _parse_frontmatter(text)
    html = md_lib.markdown(body, extensions=["extra", "codehilite", "toc"]) if HAS_MD else body
    return {"content": body, "html": html,
        "metadata": {"frontmatter": fm, "has_frontmatter": bool(fm), "char_count": len(body)},
        "structure": {"frontmatter": fm, "headings": re.findall(r"^(#{1,6})\s+(.+)$", body, re.MULTILINE)},
        "tables": [], "images": [], "errors": []}

def read_xml(path):
    content = _safe_read(path)
    tree = ET.fromstring(content)
    def _elem_to_dict(e):
        r = {"tag": e.tag, "attrib": e.attrib, "text": (e.text or "").strip()}
        children = [_elem_to_dict(c) for c in e]
        if children: r["children"] = children
        return r
    return {"content": content.decode("utf-8", errors="replace"), "metadata": {"root_tag": tree.tag, "encoding": "utf-8"},
        "structure": _elem_to_dict(tree), "tables": [], "images": [], "errors": []}

def read_json(path):
    content = _safe_read(path).decode("utf-8-sig", errors="replace")
    data = json.loads(content)
    return {"content": json.dumps(data, indent=2, ensure_ascii=False),
        "metadata": {"type": str(type(data).__name__), "size_bytes": len(content)},
        "structure": data, "tables": [], "images": [], "errors": []}

def _fallback_text(path, fmt):
    try:
        encoding = _detect_encoding(path)
        with open(path, "r", encoding=encoding) as f: content = f.read()
    except Exception:
        try: content = _safe_read(path).decode("utf-8", errors="replace")
        except Exception: content = ""
    return {"content": content, "metadata": {"fallback": True, "format": fmt}, "structure": {}, "tables": [], "images": [], "errors": [f"Required library not available for {fmt}, read as plain text"]}
