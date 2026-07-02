import io, os
from typing import Any

try:
    from docx import Document; from docx.shared import Inches, Pt, RGBColor; from docx.enum.text import WD_ALIGN_PARAGRAPH; HAS_DOCX = True
except ImportError: HAS_DOCX = False
try:
    import xlsxwriter; HAS_XLSX = True
except ImportError: HAS_XLSX = False
try:
    from pptx import Presentation; from pptx.util import Inches as PptInches, Pt as PptPt; HAS_PPTX = True
except ImportError: HAS_PPTX = False
try:
    from weasyprint import HTML; HAS_WEASYPRINT = True
except (ImportError, OSError): HAS_WEASYPRINT = False
try:
    import jinja2; HAS_JINJA2 = True
except ImportError: HAS_JINJA2 = False

TEMPLATES = {
    "propuesta_tecnica": {"title": "Propuesta Técnica", "sections": ["resumen_ejecutivo", "alcance", "metodologia", "cronograma", "equipo", "presupuesto"]},
    "informe_requerimientos": {"title": "Informe de Requerimientos", "sections": ["introduccion", "requerimientos_funcionales", "requerimientos_no_funcionales", "restricciones"]},
    "timesheet": {"title": "Registro de Horas", "sections": ["empleado", "proyecto", "horas", "descripcion"]},
    "diagrama_arquitectura": {"title": "Diagrama de Arquitectura", "sections": ["vista_general", "componentes", "interacciones", "tecnologias"]},
    "acta_reunion": {"title": "Acta de Reunión", "sections": ["fecha", "asistentes", "agenda", "discusion", "acuerdos", "proximos_pasos"]},
    "presentacion_proyecto": {"title": "Presentación del Proyecto", "sections": ["titulo", "problema", "solucion", "beneficios", "roadmap"]},
    "default": {"title": "Documento", "sections": ["contenido"]},
}

def generate_docx(content, template="default"):
    if not HAS_DOCX: raise ImportError("python-docx is required")
    doc = Document(); tmpl = TEMPLATES.get(template, TEMPLATES["default"])
    title = content.get("title", tmpl["title"]); body = content.get("text", content.get("content", ""))
    doc.add_heading(title, level=1)
    for section in tmpl["sections"]:
        doc.add_heading(section.replace("_", " ").title(), level=2)
        sc = content.get(section, content.get("text", ""))
        if isinstance(sc, list):
            for item in sc: doc.add_paragraph(str(item), style="List Bullet")
        else: doc.add_paragraph(str(sc))
    for table_data in content.get("tables", []):
        rows = table_data.get("rows", [])
        if rows:
            table = doc.add_table(rows=len(rows), cols=len(rows[0])); table.style = "Light Grid Accent 1"
            for i, row_data in enumerate(rows):
                for j, cell_data in enumerate(row_data): table.rows[i].cells[j].text = str(cell_data)
    buf = io.BytesIO(); doc.save(buf); buf.seek(0); return buf.read()

def generate_xlsx(data, template="default"):
    if not HAS_XLSX: raise ImportError("xlsxwriter is required")
    buf = io.BytesIO(); wb = xlsxwriter.Workbook(buf)
    header_fmt = wb.add_format({"bold": True, "bg_color": "#4472C4", "font_color": "white"})
    if "sheets" in data:
        for s in data["sheets"]:
            ws = wb.add_worksheet(s.get("name", "Sheet1"))
            for i, row in enumerate(s.get("rows", [])):
                for j, val in enumerate(row): ws.write(i, j, val)
    else:
        ws = wb.add_worksheet("Datos"); headers = data.get("headers", []); rows = data.get("rows", [])
        for j, h in enumerate(headers): ws.write(0, j, h, header_fmt)
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row): ws.write(i, j, val)
    wb.close(); buf.seek(0); return buf.read()

def generate_pptx(content, template="default"):
    if not HAS_PPTX: raise ImportError("python-pptx is required")
    prs = Presentation(); tmpl = TEMPLATES.get(template, TEMPLATES["default"])
    title = content.get("title", tmpl["title"])
    slide = prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text = title
    body = content.get("text", content.get("content", ""))
    if body:
        slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = "Contenido"
        if len(slide.placeholders) > 1: slide.placeholders[1].text = str(body)
    for section in tmpl["sections"]:
        sc = content.get(section, "")
        if sc:
            slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = section.replace("_", " ").title()
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "\n".join(f"\u2022 {item}" for item in sc) if isinstance(sc, list) else str(sc)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf.read()

def generate_pdf(content, options=None):
    if not HAS_WEASYPRINT: raise ImportError("WeasyPrint is required")
    options = options or {}
    title = options.get("title", "Documento"); css = options.get("css", "")
    html_str = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{title}</title>
<style>body{{font-family:'Segoe UI',Arial,sans-serif;margin:2cm;font-size:11pt;line-height:1.6}}
h1{{color:#1a365d;border-bottom:2px solid #2b6cb0;padding-bottom:8px}}
h2{{color:#2b6cb0;margin-top:24px}} p{{margin:8px 0}}
table{{border-collapse:collapse;width:100%;margin:16px 0}}
th,td{{border:1px solid #cbd5e0;padding:8px;text-align:left}}
th{{background-color:#edf2f7;font-weight:bold}}
code{{background:#f7fafc;padding:2px 6px;border-radius:3px;font-size:0.9em}}
pre{{background:#f7fafc;padding:16px;border-radius:6px;overflow-x:auto}}
{css}</style></head><body>{content}</body></html>"""
    buf = io.BytesIO(); HTML(string=html_str).write_pdf(buf); buf.seek(0); return buf.read()

def generate_markdown(content):
    title = content.get("title", "Documento"); body = content.get("text", content.get("content", ""))
    md = f"# {title}\n\n"
    for section in content.get("sections", []):
        md += f"## {section.get('title', 'Sección')}\n\n{section.get('content', '')}\n\n"
        for item in section.get("items", []): md += f"- {item}\n"
        md += "\n"
    for table in content.get("tables", []):
        if table.get("caption"): md += f"**{table['caption']}**\n\n"
        rows = table.get("rows", [])
        if rows:
            md += "| " + " | ".join(str(h) for h in rows[0]) + " |\n"
            md += "| " + " | ".join("---" for _ in rows[0]) + " |\n"
            for row in rows[1:]: md += "| " + " | ".join(str(c) for c in row) + " |\n"
        md += "\n"
    if body and not content.get("sections"): md += f"{body}\n\n"
    for lst in content.get("lists", []):
        for item in lst: md += f"- {item}\n"
        md += "\n"
    return md
