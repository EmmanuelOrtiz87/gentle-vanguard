from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette - Executive Board Style
NAVY = RGBColor(20, 40, 80)
DARK_GRAY = RGBColor(45, 45, 50)
MEDIUM_GRAY = RGBColor(80, 80, 90)
LIGHT_GRAY = RGBColor(240, 242, 245)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(220, 60, 40)  # Red for CTAs
SUCCESS = RGBColor(0, 140, 100)  # Green for positive metrics
ACCENT_BLUE = RGBColor(0, 120, 200)

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(180, 190, 200)
        p.font.name = "Calibri"
    
    return slide

def add_section_slide(title, section_number=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY
    shape.line.fill.background()
    
    if section_number:
        num_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_number
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.font.name = "Calibri"
    
    title_box = slide.shapes.add_textbox(Inches(1) if not section_number else Inches(3.5), Inches(2.5), Inches(9.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    return slide

def add_kpi_slide(title, kpis):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    num_kpis = len(kpis)
    box_width = 3.5 if num_kpis <= 3 else (12.333 / num_kpis) - 0.3
    spacing = 0.5 if num_kpis <= 3 else 0.2
    total_width = num_kpis * box_width + (num_kpis - 1) * spacing
    start_x = (13.333 - total_width) / 2
    
    for i, kpi in enumerate(kpis):
        number = kpi[0]
        label = kpi[1]
        is_accent = kpi[2] if len(kpi) > 2 else False
        
        x = start_x + i * (box_width + spacing)
        
        # KPI Box
        kpi_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2), Inches(box_width), Inches(2.8))
        kpi_box.fill.solid()
        kpi_box.fill.fore_color.rgb = ACCENT if is_accent else LIGHT_GRAY
        kpi_box.line.fill.background()
        
        # Number
        num_text = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(box_width), Inches(1.2))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = WHITE if is_accent else NAVY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_text = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.5), Inches(box_width - 0.2), Inches(1.2))
        tf = label_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE if is_accent else DARK_GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_impact_slide(title, main_point, supporting_points=None, call_to_action=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    # Main impact point - large and bold
    main_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.5), Inches(2))
    tf = main_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_point
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"
    
    # Supporting points as small cards
    if supporting_points:
        y = 4
        for point in supporting_points:
            # Small indicator
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.15), Inches(0.15), Inches(0.15))
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT_BLUE
            dot.line.fill.background()
            
            point_text = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(11), Inches(0.6))
            tf = point_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = MEDIUM_GRAY
            p.font.name = "Calibri"
            y += 0.7
    
    # Call to action box
    if call_to_action:
        cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(6.2), Inches(4.8), Inches(0.9))
        cta_box.fill.solid()
        cta_box.fill.fore_color.rgb = ACCENT
        cta_box.line.fill.background()
        
        cta_text = slide.shapes.add_textbox(Inches(8), Inches(6.35), Inches(4.8), Inches(0.6))
        tf = cta_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = call_to_action
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_process_slide(title, steps, flow_direction="horizontal"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    num_steps = len(steps)
    box_width = 2.2
    spacing = 0.4
    total_width = num_steps * box_width + (num_steps - 1) * spacing
    start_x = (13.333 - total_width) / 2
    y = 2.5
    
    for i, step_data in enumerate(steps):
        x = start_x + i * (box_width + spacing)
        
        # Handle both formats: (num, title, desc) or (name, price, desc)
        if len(step_data) == 3:
            if isinstance(step_data[0], int):
                step_num, step_title, step_desc = step_data
                show_number = True
            else:
                step_title, step_desc, show_number = step_data[0], step_data[1], False
                step_num = ""
        else:
            step_title, step_desc = step_data[0], step_data[1]
            show_number = False
            step_num = ""
        
        # Step box
        step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_width), Inches(3.2))
        step_box.fill.solid()
        # Highlight WF as first option
        if i == 0:
            step_box.fill.fore_color.rgb = SUCCESS
        else:
            step_box.fill.fore_color.rgb = LIGHT_GRAY
        step_box.line.fill.background()
        
        # Number or Title circle
        if show_number:
            num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.8), Inches(y + 0.2), Inches(0.6), Inches(0.6))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = ACCENT_BLUE
            num_circle.line.fill.background()
            
            num_text = slide.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.25), Inches(0.6), Inches(0.5))
            tf = num_text.text_frame
            p = tf.paragraphs[0]
            p.text = str(step_num)
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        else:
            # Title as label
            title_label = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(box_width), Inches(0.5))
            tf = title_label.text_frame
            p = tf.paragraphs[0]
            p.text = step_title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WHITE if i == 0 else NAVY
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_text = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.8), Inches(box_width - 0.2), Inches(2.2))
        tf = desc_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE if i == 0 else MEDIUM_GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow between steps
        if i < num_steps - 1:
            arrow_x = x + box_width + 0.05
            arrow = slide.shapes.add_textbox(Inches(arrow_x), Inches(y + 1.3), Inches(0.3), Inches(0.5))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(24)
            p.font.color.rgb = ACCENT_BLUE
            p.font.name = "Calibri"
    
    return slide

def add_comparison_slide(title, before, after):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    # Before box
    before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5))
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RGBColor(250, 235, 235)
    before_box.line.color.rgb = RGBColor(200, 100, 100)
    before_box.line.width = Pt(2)
    
    before_label = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(5.4), Inches(0.6))
    tf = before_label.text_frame
    p = tf.paragraphs[0]
    p.text = "ANTES"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 60, 60)
    p.font.name = "Calibri"
    
    before_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(5.4), Inches(3.8))
    tf = before_content.text_frame
    tf.word_wrap = True
    for i, point in enumerate(before):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✗ " + point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(100, 60, 60)
        p.font.name = "Calibri"
        p.space_after = Pt(12)
    
    # Arrow
    arrow = slide.shapes.add_textbox(Inches(6.5), Inches(3.8), Inches(0.5), Inches(1))
    tf = arrow.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # After box
    after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = RGBColor(235, 250, 240)
    after_box.line.color.rgb = SUCCESS
    after_box.line.width = Pt(2)
    
    after_label = slide.shapes.add_textbox(Inches(7.2), Inches(2), Inches(5.4), Inches(0.6))
    tf = after_label.text_frame
    p = tf.paragraphs[0]
    p.text = "DESPUÉS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.font.name = "Calibri"
    
    after_content = slide.shapes.add_textbox(Inches(7.2), Inches(2.7), Inches(5.4), Inches(3.8))
    tf = after_content.text_frame
    tf.word_wrap = True
    for i, point in enumerate(after):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✓ " + point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 100, 60)
        p.font.name = "Calibri"
        p.space_after = Pt(12)
    
    return slide

def add_timeline_slide(title, phases):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    
    y = 1.8
    colors = [SUCCESS, ACCENT_BLUE, RGBColor(180, 140, 0), ACCENT]
    
    for i, (phase, duration, deliverable, is_key) in enumerate(phases):
        # Timeline dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.3), Inches(0.4), Inches(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors[i % len(colors)]
        dot.line.fill.background()
        
        # Timeline line
        if i < len(phases) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y + 0.7), Inches(0.1), Inches(1))
            line.fill.solid()
            line.fill.fore_color.rgb = LIGHT_GRAY
            line.line.fill.background()
        
        # Phase card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y), Inches(11.3), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY if not is_key else RGBColor(255, 245, 235)
        card.line.fill.background()
        
        # Phase title
        phase_text = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.1), Inches(6), Inches(0.5))
        tf = phase_text.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY if not is_key else RGBColor(180, 100, 40)
        p.font.name = "Calibri"
        
        # Duration
        dur_text = slide.shapes.add_textbox(Inches(10), Inches(y + 0.1), Inches(2.5), Inches(0.5))
        tf = dur_text.text_frame
        p = tf.paragraphs[0]
        p.text = duration
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_BLUE if not is_key else ACCENT
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.RIGHT
        
        # Deliverable
        del_text = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.6), Inches(10.8), Inches(0.7))
        tf = del_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = deliverable
        p.font.size = Pt(14)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = "Calibri"
        
        y += 1.6
    
    return slide

def add_cta_slide(title, action_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # Action items
    y = 2.5
    for i, (num, action) in enumerate(action_items):
        # Number box
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y), Inches(0.7), Inches(0.7))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = ACCENT
        num_box.line.fill.background()
        
        num_text = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.1), Inches(0.7), Inches(0.5))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        
        # Action text
        action_text = slide.shapes.add_textbox(Inches(2.4), Inches(y + 0.1), Inches(9.5), Inches(0.6))
        tf = action_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = action
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        
        y += 1.1
    
    # Bottom accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(6.8), Inches(3.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    return slide

# =====================================================
# SLIDES
# =====================================================

# Slide 1: Title
add_title_slide(
    "Workspace Foundation",
    "Estandarización y Aceleración del Desarrollo con Inteligencia Artificial"
)

# Slide 2: Executive Summary
add_kpi_slide(
    "Resumen Ejecutivo",
    [
        ("50%", "Reducción en tiempo de configuración"),
        ("30%", "Aumento en productividad"),
        ("0", "Costo de licencia"),
        ("1-click", "Onboarding", True)
    ]
)

# Slide 3: The Problem
add_impact_slide(
    "El Problema",
    "Cada proyecto comienza igual: configuración manual, herramientas dispersas, tiempo perdido.",
    [
        "Configuración promedio: 2-4 horas por desarrollador",
        "Herramientas AI sin governance ni estándares",
        "Inconsistencia entre equipos y proyectos",
        "Sin visibilidad del uso de herramientas AI",
        "Reinventar la rueda en cada nuevo proyecto"
    ]
)

# Slide 4: Our Solution
add_impact_slide(
    "La Solución",
    "Un sistema unificado que configura, estandariza y mide el desarrollo con IA.",
    [
        "Bootstrap automático: un comando, todo listo",
        "Plantillas inteligentes para cada tipo de proyecto",
        "Integración nativa con Claude, OpenCode, GGA",
        "Sistema de auditoría integrado",
        "Documentación y estándares automatizados"
    ],
    call_to_action="Ver arquitectura →"
)

# Slide 5: Architecture Overview
add_process_slide(
    "Arquitectura del Sistema",
    [
        (1, "Bootstrap", "Un script configura todo automáticamente"),
        (2, "Templates", "Estructuras predefinidas por tipo de proyecto"),
        (3, "Skills", "Patrones reutilizables de IA"),
        (4, "Audit", "Tracking automático de actividad"),
        (5, "Docs", "Documentación integrada")
    ]
)

# Slide 6: Before vs After
add_comparison_slide(
    "Impacto Medible",
    [
        "2-4 horas configurando entorno",
        "Documentación inconsistente",
        "Herramientas instaladas manualmente",
        "Sin estándares de código",
        "Zero visibilidad de AI usage"
    ],
    [
        "5 minutos listo para trabajar",
        "Documentación automática",
        "Un comando instala todo",
        "Estándares integrados desde inicio",
        "Métricas y reportes automáticos"
    ]
)

# Slide 7: Key Features
add_impact_slide(
    "Capacidades Clave",
    "Zero-config. Máxima productividad.",
    [
        "Multi-provider AI: Claude, OpenAI, Gemini compatibles",
        "GGA: Code review automatizado con IA",
        "Gentle-AI: Asistente contextual integrado",
        "Skills reutilizables para patrones comunes",
        "Generación de código, tests, docs con un comando"
    ]
)

# Slide 8: Audit & Metrics
add_kpi_slide(
    "Auditoría y Métricas",
    [
        ("100%", "Transparencia"),
        ("7 días", "Retención histórica"),
        ("$0", "Costo adicional"),
        ("Git", "Versionado", True)
    ]
)

# Slide 9: Implementation Timeline
add_timeline_slide(
    "Roadmap de Adopción",
    [
        ("Semana 1-2: Piloto", "2 semanas", "Implementar en 2-3 proyectos seleccionados", False),
        ("Semana 3-4: Feedback", "2 semanas", "Ajustar según retroalimentación del equipo", False),
        ("Mes 2: Rollout", "1 mes", "Extender a todos los equipos", True),
        ("Mes 3+: Optimización", "Ongoing", "Mejora continua basada en métricas", False)
    ]
)

# Slide 10: Prerequisites
add_impact_slide(
    "Requisitos Previos",
    "Mínimo. Máximo impacto.",
    [
        "Git instalado (cualquier versión)",
        "PowerShell 5.0+ o Bash (Linux/Mac)",
        "Acceso a internet para descargar herramientas",
        "Credenciales de API (Claude, OpenAI) opcionales",
        "Repositorio Git (existente o nuevo)"
    ]
)

# Slide 11: Investment
add_kpi_slide(
    "Inversión Requerida",
    [
        ("$0", "Licencia"),
        ("$0", "Infraestructura"),
        ("2 sem", "Piloto", True),
        ("$0", "Training")
    ]
)

# Slide 12: ROI
add_impact_slide(
    "Retorno de Inversión",
    "Tiempo recuperado en el primer mes.",
    [
        "Piloto: 2 semanas de implementación",
        "Ahorro: 2 horas × 10 devs × 20 días = 400 horas/mes",
        "Valor: $400/hr × 400 hrs = $160,000/mes",
        "ROI: Inmediato desde el día 1",
        "Escala linealmente con el equipo"
    ],
    call_to_action="Ver análisis completo →"
)

# Slide 13: Team Adoption
add_impact_slide(
    "Adopción por Equipo",
    "El desarrollador promedio recupera 1 hora/día.",
    [
        "Onboarding de nuevos devs: de 2 días a 2 horas",
        "Consistencia entre proyectos: mismo baseline para todos",
        "Reducción de bugs: estándares integrados",
        "Menos contexto switching: herramientas unificadas",
        "Satisfacción del equipo: menos fricción, más código"
    ],
    call_to_action="Ver plan de training →"
)

# Slide 14: Training Plan
add_impact_slide(
    "Capacitación del Equipo",
    "16 horas de training. Documentación lista.",
    [
        "Technical Onboarding Guide creado (120+ páginas)",
        "Quick Reference Card para cada dev",
        "4 semanas de training estructurado",
        "Champions por equipo para support",
        "El equipo puede iniciar en 1 semana"
    ],
    call_to_action="Ver TECHNICAL-ONBOARDING.md →"
)

# Slide 15: Success Metrics
add_kpi_slide(
    "Métricas de Éxito",
    [
        ("+25%", "Productividad"),
        ("-50%", "Tiempo setup"),
        (">80%", "Adoption"),
        (">35", "NPS equipo", True)
    ]
)

# Slide 16: Risks & Mitigation
add_impact_slide(
    "Riesgos y Mitigaciones",
    "Todos los riesgos mapeados con acciones concretas.",
    [
        "Dependencia IA: Código siempre revisado por humano",
        "Info sensible: Políticas de secure prompting",
        "Vendor lock-in: Multi-provider (Claude + OpenAI)",
        "Costos descontrolados: Audit system + alerts",
        "Resistencia del equipo: Champions + communication"
    ],
    call_to_action="Plan completo en BOARD-SUPPLEMENT.md"
)

# Slide 17: Competition Comparison
add_process_slide(
    "Comparativa con Alternativas",
    [
        ("WF", "Gratis", "Full platform"),
        ("Copilot", "$19/mo", "Solo MS"),
        ("Amazon Q", "$19/mo", "Solo AWS"),
        ("Cursor", "$20/mo", "No governance"),
        ("Nada", "$0", "Sin estándares")
    ]
)

# Slide 18: Why WF
add_kpi_slide(
    "Por Qué Workspace Foundation",
    [
        ("$0", "Licencia", True),
        ("Multi", "Vendor", False),
        ("Audit", "Built-in", False),
        ("Open", "Source", False)
    ]
)

# Slide 19: Call to Action
add_cta_slide(
    "Llamada a la Acción",
    [
        (1, "Aprobar piloto: 2-3 proyectos, 4 semanas"),
        (2, "Designar champion técnico por equipo"),
        (3, "Iniciar training (material listo)"),
        (4, "Primeras métricas en 30 días")
    ]
)

# Slide 20: Thank You
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "¿Preguntas?"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
tf = subtitle_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Workspace Foundation\nEstandarizar. Acelerar. Medir."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(180, 190, 200)
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.5), Inches(2.333), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# CTA reminder
cta_reminder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(6.2), Inches(7.333), Inches(0.9))
cta_reminder.fill.solid()
cta_reminder.fill.fore_color.rgb = ACCENT
cta_reminder.line.fill.background()

cta_text = slide.shapes.add_textbox(Inches(3), Inches(6.35), Inches(7.333), Inches(0.6))
tf = cta_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Podemos iniciar en 1 semana"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "C:\\Workspace_local\\workspace-foundation\\docs\\WORKSPACE-FOUNDATION-CONSEJO.pptx"
prs.save(output_path)
print(f"Presentación guardada: {output_path}")
